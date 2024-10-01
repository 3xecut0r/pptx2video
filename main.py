import argparse
import logging
from gtts import gTTS
from pptx import Presentation
from moviepy.editor import *
import os
from typing import List, Optional

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


class SlideExtractor:
    def __init__(self, pptx_file: str):
        self.pptx_file = pptx_file
        self.slides: List[str] = []
        self.notes: List[str] = []

    def extract_slides_and_notes(self) -> None:
        presentation = Presentation(self.pptx_file)
        for i, slide in enumerate(presentation.slides):
            image_filename = self._save_slide_as_image(slide, i)
            if image_filename:
                self.slides.append(image_filename)
            if slide.has_notes_slide:
                note_text = slide.notes_slide.notes_text_frame.text
                self.notes.append(note_text)

    @staticmethod
    def _save_slide_as_image(slide, index: int) -> Optional[str]:
        for shape in slide.shapes:
            if hasattr(shape, 'image') and shape.shape_type == 13:
                image = shape.image
                image_bytes = image.blob
                image_filename = f'slide_{index}.png'
                with open(image_filename, 'wb') as f:
                    f.write(image_bytes)
                return image_filename
        return None


class AudioGenerator:
    def __init__(self, notes: List[str], lang: str = 'en'):
        self.notes = notes
        self.lang = lang
        self.audio_files: List[str] = []

    def generate_audio(self) -> None:
        for i, note in enumerate(self.notes):
            try:
                audio_filename = f'note_{i}.mp3'
                logging.info(f'Creating audio file for note {i}')
                tts = gTTS(note, lang=self.lang)
                tts.save(audio_filename)
                self.audio_files.append(audio_filename)
            except Exception as e:
                logging.error(f'Error creating audio file for note {i}: {e}')


class VideoGenerator:
    def __init__(self, slides: List[str], audio_files: List[str]):
        self.slides = slides
        self.audio_files = audio_files

    def create_video(self, output_file: str = 'presentation_video.mp4') -> None:
        video_clips = []
        for i, slide_image in enumerate(self.slides):
            audio_filename = self.audio_files[i]
            if os.path.exists(audio_filename):
                image_clip = ImageClip(slide_image)
                audio_clip = AudioFileClip(audio_filename)
                image_clip = image_clip.set_duration(audio_clip.duration)
                image_clip = image_clip.set_audio(audio_clip)
                video_clips.append(image_clip)
            else:
                logging.warning(f'Skipping slide {i}, as audio file {audio_filename} not found.')

        if video_clips:
            final_video = concatenate_videoclips(video_clips, method='compose')
            final_video.write_videofile(output_file, fps=24)

    def cleanup(self) -> None:
        for slide_image in self.slides:
            os.remove(slide_image)
        for audio_file in self.audio_files:
            if os.path.exists(audio_file):
                os.remove(audio_file)


class PresentationToVideoConverter:
    def __init__(self, pptx_file: str, lang: str = 'en'):
        self.pptx_file = pptx_file
        self.lang = lang

    def convert(self, output_file: str) -> None:
        slide_extractor = SlideExtractor(self.pptx_file)
        slide_extractor.extract_slides_and_notes()

        audio_generator = AudioGenerator(slide_extractor.notes, self.lang)
        audio_generator.generate_audio()

        video_generator = VideoGenerator(slide_extractor.slides, audio_generator.audio_files)
        video_generator.create_video(output_file)

        video_generator.cleanup()


def main():
    parser = argparse.ArgumentParser(description='Convert PowerPoint presentation to a video with audio narration.')
    parser.add_argument('--pptx', required=True, help='Path to the PowerPoint (.pptx) file')
    parser.add_argument('--output', default='presentation_video.mp4', help='Name of the output video file (.mp4)')
    parser.add_argument('--lang', default='en', help='Language for the narration (e.g., en, de, fr)')

    args = parser.parse_args()

    converter = PresentationToVideoConverter(pptx_file=args.pptx, lang=args.lang)
    converter.convert(output_file=args.output)


if __name__ == '__main__':
    main()
